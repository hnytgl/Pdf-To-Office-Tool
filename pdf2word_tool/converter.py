from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

OUTPUT_SUFFIXES = {
    "word": ".docx",
    "docx": ".docx",
    "ppt": ".pptx",
    "pptx": ".pptx",
    "excel": ".xlsx",
    "xlsx": ".xlsx",
}


@dataclass(frozen=True)
class ConversionJob:
    source: Path
    target: Path
    output_format: str


@dataclass(frozen=True)
class ConversionResult:
    source: Path
    target: Path
    status: str
    message: str = ""


def normalize_format(output_format: str) -> str:
    value = output_format.lower()
    if value in {"docx", "word"}:
        return "word"
    if value in {"pptx", "ppt"}:
        return "ppt"
    if value in {"xlsx", "excel"}:
        return "excel"
    raise ValueError(f"不支持的输出格式：{output_format}")


def suffix_for_format(output_format: str) -> str:
    return OUTPUT_SUFFIXES[normalize_format(output_format)]


def discover_jobs(
    input_path: Path,
    output_path: Path | None = None,
    recursive: bool = False,
    output_format: str = "word",
) -> list[ConversionJob]:
    input_path = input_path.expanduser().resolve()
    if output_path is not None:
        output_path = output_path.expanduser().resolve()
    normalized_format = normalize_format(output_format)
    suffix = suffix_for_format(normalized_format)

    if input_path.is_file():
        if input_path.suffix.lower() != ".pdf":
            raise ValueError(f"输入文件不是 PDF：{input_path}")
        target = output_path or input_path.with_suffix(suffix)
        if target.suffix.lower() != suffix:
            target = target / input_path.with_suffix(suffix).name
        return [ConversionJob(input_path, target, normalized_format)]

    if not input_path.is_dir():
        raise FileNotFoundError(f"找不到输入路径：{input_path}")

    pattern = "**/*.pdf" if recursive else "*.pdf"
    pdf_files = sorted(input_path.glob(pattern))
    target_root = output_path or input_path
    jobs = []
    for pdf_file in pdf_files:
        relative = pdf_file.relative_to(input_path)
        target = target_root / relative.with_suffix(suffix)
        jobs.append(ConversionJob(pdf_file, target, normalized_format))
    return jobs


def convert_jobs(
    jobs: Iterable[ConversionJob],
    overwrite: bool = False,
    start_page: int | None = None,
    end_page: int | None = None,
) -> list[ConversionResult]:
    results: list[ConversionResult] = []
    for job in jobs:
        results.append(convert_one(job, overwrite=overwrite, start_page=start_page, end_page=end_page))
    return results


def convert_one(
    job: ConversionJob,
    overwrite: bool = False,
    start_page: int | None = None,
    end_page: int | None = None,
) -> ConversionResult:
    if job.target.exists() and not overwrite:
        return ConversionResult(job.source, job.target, "skipped", "目标文件已存在，使用 --overwrite 可覆盖")

    job.target.parent.mkdir(parents=True, exist_ok=True)
    if job.output_format == "word":
        return convert_word(job, start_page=start_page, end_page=end_page)
    if job.output_format == "ppt":
        return convert_ppt(job, start_page=start_page, end_page=end_page)
    if job.output_format == "excel":
        return convert_excel(job, start_page=start_page, end_page=end_page)
    return ConversionResult(job.source, job.target, "failed", f"不支持的输出格式：{job.output_format}")


def convert_word(job: ConversionJob, start_page: int | None = None, end_page: int | None = None) -> ConversionResult:
    try:
        from pdf2docx import Converter
    except Exception as exc:
        return ConversionResult(job.source, job.target, "failed", f"未安装 pdf2docx：{exc}")

    start = (start_page - 1) if start_page is not None else 0
    end = end_page if end_page is not None else None

    try:
        converter = Converter(str(job.source))
        try:
            converter.convert(str(job.target), start=start, end=end)
        finally:
            converter.close()
    except Exception as exc:
        return ConversionResult(job.source, job.target, "failed", str(exc))

    return ConversionResult(job.source, job.target, "converted")


def convert_ppt(job: ConversionJob, start_page: int | None = None, end_page: int | None = None) -> ConversionResult:
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        return ConversionResult(job.source, job.target, "failed", f"未安装 PPT 转换依赖：{exc}")

    try:
        doc = fitz.open(job.source)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        first_page = start_page or 1
        last_page = end_page or doc.page_count
        temp_images: list[Path] = []

        for page_number in range(first_page, min(last_page, doc.page_count) + 1):
            page = doc.load_page(page_number - 1)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = job.target.parent / f".{job.target.stem}_page_{page_number}.png"
            pixmap.save(image_path)
            temp_images.append(image_path)

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

        if not temp_images:
            return ConversionResult(job.source, job.target, "failed", "指定页码范围内没有可转换页面")

        prs.save(job.target)
        for image_path in temp_images:
            image_path.unlink(missing_ok=True)
        doc.close()
    except Exception as exc:
        return ConversionResult(job.source, job.target, "failed", str(exc))

    return ConversionResult(job.source, job.target, "converted")


def convert_excel(job: ConversionJob, start_page: int | None = None, end_page: int | None = None) -> ConversionResult:
    try:
        import openpyxl
        import pdfplumber
    except Exception as exc:
        return ConversionResult(job.source, job.target, "failed", f"未安装 Excel 转换依赖：{exc}")

    try:
        workbook = openpyxl.Workbook()
        default_sheet = workbook.active
        workbook.remove(default_sheet)
        with pdfplumber.open(job.source) as pdf:
            first_page = start_page or 1
            last_page = end_page or len(pdf.pages)
            selected_pages = pdf.pages[first_page - 1 : min(last_page, len(pdf.pages))]
            if not selected_pages:
                return ConversionResult(job.source, job.target, "failed", "指定页码范围内没有可转换页面")

            for index, page in enumerate(selected_pages, start=first_page):
                worksheet = workbook.create_sheet(title=f"第{index}页")
                tables = page.extract_tables()
                row_index = 1
                if tables:
                    for table in tables:
                        for row in table:
                            for col_index, value in enumerate(row, start=1):
                                worksheet.cell(row=row_index, column=col_index, value=value)
                            row_index += 1
                        row_index += 1
                else:
                    text = page.extract_text() or ""
                    for line in text.splitlines():
                        worksheet.cell(row=row_index, column=1, value=line)
                        row_index += 1
                    if not text:
                        worksheet.cell(row=1, column=1, value="该页未提取到文本或表格")

        workbook.save(job.target)
    except Exception as exc:
        return ConversionResult(job.source, job.target, "failed", str(exc))

    return ConversionResult(job.source, job.target, "converted")
